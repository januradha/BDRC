#!/usr/bin/python
# print "Content-Type: text/html; charset=utf-8\n"

import os
import sys
import math
import cgi
import datetime
import tempfile
import collections

from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import ChartData
import pptx

import model
import get_scores
import brand
import db
from bdrcshared.demense import Demense
from division import Division
import venue
from reporting_base import ReportBase

from common_report import CommonCall
import common


DEM = 'DEM'
DIV = 'DIV'
QID = 'QID'

inputFile = './templates/BDRC_2016_Input.pptx'
prs = Presentation(inputFile)

## temp file is used to keep generated output file from inputFile
temp = tempfile.NamedTemporaryFile(suffix='.pptx', prefix='BDRC_2016_Output_', dir='/tmp')
outputFile = temp.name

## desttemp file is used to rearrange slide for outputFile file.
## This will call after Output file generated and saved
desttemp = tempfile.NamedTemporaryFile(suffix='.pptx', prefix='BDRC_2016_Output_', dir='/tmp')
desttempFile = desttemp.name


class ReportManager:

    def __init__(self, dbConnection):

        if dbConnection:
            self.d = dbConnection
        else:
            import db
            self.d = db.connect()

    def generatePPT(self,
                    brand_id,
                    brand_name,
                    demense_id,
                    demense_name,
                    score_month,
                    score_month_value,
                    score_year,
                    division_id,
                    selected_index):
        """
        On click of Generate PPTX button this method will execute 
        and for each slide.

        @brand_id  :  (int) Brand ID  - 225 for Hilton
        @brand_name : (String) Brand name - "Hilton Worldwide" 
        @demense_id : (int) Region ID
        @demense_name : (String) Demense name
        @score_month :  (int) 3 for March
        @score_month_value : (String) March
        @score_year : (int) 2016
        @division_id : (int) Division ID
        """
        selected_slide_list,remove_slide_index = self.preConditionReportGenerate(selected_index)

        no_of_extra_slide = self.reportGenerate(brand_id, brand_name, demense_id, demense_name, score_month, score_month_value, score_year, division_id,selected_slide_list)

        self.postConditionReportGenerate(no_of_extra_slide,selected_index,selected_slide_list,remove_slide_index)

    def preConditionReportGenerate(self,
                                   selected_index):
        """
        This will Check user's slide selection list
        Based on that it will generate a list for slide design
        For e.g 
        Total range(1,13)
        deselect 5,9
        selected_slide_list = [1,2,3,4,6,7,8,10,11,12]
        remove_slide_index = [5,9]
        """
        remove_slide_index = []
        prs_max_length = range((prs.slides.__len__()))

        if selected_index:
            selected_slide_list = []
            selected_slide_list.append(0)
            selected_index = [int(x) for x in selected_index.split(',')]
            selected_slide_list.extend(selected_index)
            selected_slide_list.append(11)
            remove_slide_index = list(set(prs_max_length).difference(selected_slide_list))
        else:
            selected_slide_list = prs_max_length

        return selected_slide_list,remove_slide_index

    def reportGenerate(self,
                       brand_id,
                       brand_name,
                       demense_id,
                       demense_name,
                       score_month,
                       score_month_value,
                       score_year,
                       division_id,
                       selected_slide_list):
        """
        Main Method To Generate Report
        """
        score_type = '1'
        answer_value = '!'
        id_type = 0
        no_of_extra_slide = None

        gpptx = GeneratePPT(self.d,
                            brand_id=int(brand_id),
                            demense_name=demense_name,
                            score_month=int(score_month),
                            score_month_value=score_month_value,
                            score_year=int(score_year),
                            brand_name=brand_name,
                            questionnaire_id=int(GeneratePPT.M_OVERALL_INC_QUICKCHECK),
                            id_type=id_type,
                            answer_value=answer_value,
                            demense_id=int(demense_id),
                            division_id=int(division_id),
                            score_type=score_type,
                            grouping_level=2,
                            )

        # selected_slide_list = [0, 5, 6]
        for source_index in selected_slide_list:
            mapped_index = source_index + 1
            ## calls each slides method to update slide information
            output = getattr(gpptx, 'mappedSlide%d' %
                             mapped_index)(source_index=source_index)
            ## Currently for slide 5 and 10 we need all record to display.
            ## To adjust extra nos of slide in proper order
            if source_index in (5, 10):
                no_of_extra_slide = output
        prs.save(outputFile)

        return no_of_extra_slide

    def postConditionReportGenerate(self,
                                    no_of_extra_slide,
                                    selected_index,
                                    selected_slide_list,
                                    remove_slide_index):
        """
        * update_slides will update dynamic context
        * move_slides will move extra slide with appopiate parent slide's position
        * rearrange_slides slides is required after removing wantted slides
        """

        #update_slides should be called on the required slide list
        self.update_slides(score_month=int(score_month),
                           score_year=int(score_year),
                           brand_name=brand_name,
                           demense_name=demense_name)

        ## Manage Extra no of slides with appopiate parent slide's position
        if no_of_extra_slide and isinstance(no_of_extra_slide, dict):
            self.move_slides(no_of_extra_slide)
            ## Manage remove_slide_index as per user deselect
            if selected_index and [item in selected_slide_list for item in no_of_extra_slide.keys()]:
                extra = no_of_extra_slide[item] - 1
                for i, value in enumerate(remove_slide_index):
                    if value > item:
                        remove_slide_index.remove(value)
                        remove_slide_index.insert(i, value + extra)

        ## Manage to rearrange slide index
        ## rearrange_slides slides is required after removing wantted slides
        if selected_index:
            self.rearrange_slides(selected_slide_list, no_of_extra_slide, remove_slide_index)

    def move_slide(self, presentation, old_index, new_index):
        """
        This method is used to adjust extra nos of slide by 
        moving slides order in proper order
        """
        xml_slides = presentation.slides._sldIdLst  ## pylint: disable=W0212
        slides = list(xml_slides)
        xml_slides.remove(slides[old_index])
        xml_slides.insert(new_index, slides[old_index])

    def delete_slide(self, presentation,  index):
        """
        This method is used to adjust extra nos of slide by 
        deleting unwanted slides in proper order
        """
        xml_slides = presentation.slides._sldIdLst  ## pylint: disable=W0212
        slides = list(xml_slides)
        xml_slides.remove(slides[index])

    def move_slides(self, no_of_extra_slide):
        """
        Move Slide from Old Index to New Index
        To adjust extra nos of slide in proper order
        """
        index_map = {}
        ## Object for Input file
        inputFile = './templates/BDRC_2016_Input.pptx'
        inputFileprs = os.getcwd() + '/' + inputFile
        prs = Presentation(inputFileprs)
        max_length = prs.slides.__len__()
        outputFileprs = outputFile
        desprs = Presentation(outputFileprs)
        no_of_extra_slide.keys()
        ## Manage Slide index order sortted format
        source_index_list = no_of_extra_slide.keys()
        source_index_list.sort()

        for item, source_index in enumerate(source_index_list):
            if item != 0:
                extra_slide = no_of_extra_slide[
                    source_index_list[item - 1]] - 1
                next_key = source_index + extra_slide
            else:
                next_key = source_index
            nos_of_new_record = no_of_extra_slide[source_index]
            for i in range(1, nos_of_new_record):
                key = next_key + i
                value = max_length + i - 1
                index_map[key] = value
            max_length = max_length + no_of_extra_slide[source_index] - 1

        ## Manage to arrange slide order with extra nos of slide count and index
        for new_index, old_index in index_map.iteritems():
            self.move_slide(desprs, old_index, new_index)
        desprs.save(outputFile)

    def update_slides(self,score_month,score_year,brand_name,demense_name):
        """
        update_slides is a generic method used to update all tag values with appropiate text.
        So that we can apply all in a single pass.
        """
        pyear,pmonth = common.previous_period(year=score_year,month=score_month)
        content = prs.slides
        targets = {
            '<YYYY>':   str(score_year),
            '<MONTH>':  datetime.date(1900, score_month, 1).strftime('%B'),
            '<BRAND>':  brand_name,
            '<DEMENSE>':demense_name,
            '<Y-1>':    str(pyear),
            '<M-1>':    datetime.date(1900, pmonth, 1).strftime('%B'),
        }
        for slide in content:
            self._update_slides(slide.shapes,targets)

        self._update_slides(prs.slide_master.shapes,targets,' ')
        prs.save(outputFile)

    def _update_slides(self,shapes,targets,extra=''):
        ''' update all tags within a list of shapes'''
        for shape in shapes:
            if shape.has_text_frame:    # simple text field
                self._update_paragraph(shape.text_frame.paragraphs, targets,extra)

            if shape.has_table:         # more complex table
                for row in shape.table.rows:    # drill into table rows
                    self._update_paragraph( row.cells[0].text_frame.paragraphs,targets,extra)  # only replace in 1st column (row-cell=0)

    def _update_paragraph(self,paragraphs,targets,extra=''):
        '''do a search-and-replace across a list of paragraphs'''
        for paragraph in paragraphs:
            for run in paragraph.runs:  # drilling down to this level preserves font attribs
                change,txt = self._search_and_replace(run.text,targets,extra=extra)
                if change:
                    run.text = txt


    def _search_and_replace(self,original_text,targets,extra=''):
        '''do a text search-and-replace
            return a tuple of flag & changed text'''
        change = False
        txt = original_text
        for tgt,repl in targets.items():
            if txt.find(tgt) != -1:
                txt = txt.replace(tgt,repl+extra)
                change = True

        if change:
            return True,txt

        return False,original_text


    def delete_slides(self, remove_slide_index):
        """
        Move Slide from Old Index to New Index
        """
        outputFileprs = outputFile
        desprs = Presentation(outputFileprs)
        ## Manage to remove slides
        for i in reversed(remove_slide_index):
            self.delete_slide(desprs, i)
        desprs.save(outputFile)

    def tile_order(self, mapped_slide_list, extra_slide_dict):
        """
        Set Order as per new list and extra slide
        """
        for key, val in extra_slide_dict.items():
            if key in mapped_slide_list and val:
                mapped_slide_list = [
                    myval + val - 1 if (myval > key) else myval for myval in mapped_slide_list]
                my_index = mapped_slide_list.index(key)
                for pos in range(1, val):
                    mapped_slide_list.insert(my_index + pos, key + pos)
        return mapped_slide_list

    def rearrange_slides(self, new_order, no_of_extra_slide, remove_slide_index):
        """
        Re-arrange Order as per list selecetd by user on reportorder.py
        for e.g If user has selected in order 1,3,6,2,7.

        """

        ## Manage Order as per new list and extra slide
        slide_mapped_list = self.tile_order(new_order, no_of_extra_slide)

        ## Open two copies of the initial presentation, use on a s a destination.
        ## use a sepreate src copy, so we have a constant reference
        ## to slides by slide_nr, otherwise as we rearrange them their index
        ## and slide number changes.
        outputprs = Presentation(outputFile)
        src_xml_slides = outputprs.slides._sldIdLst
        srcslides = list(src_xml_slides)
        destprs = Presentation(outputFile)
        dest_xml_slides = destprs.slides._sldIdLst

        ## Total count of slides it will set as per user selection
        final_length = len(slide_mapped_list) - len(remove_slide_index)
        for i, slide_nr in enumerate(slide_mapped_list):
            ## Only need to consider valid slides
            if i >= final_length:
                break
            ## Copy slide from src to dest.
            dest_xml_slides.insert(i, srcslides[slide_nr])

        dest_length = range(len(dest_xml_slides))
        if remove_slide_index:
            remove_slide_index = dest_length[final_length:-1]
        else:
            remove_slide_index = dest_length[final_length:]

        for i in reversed(remove_slide_index):
            self.delete_slide(destprs, i)

        destprs.save(desttempFile)

class Colours:
    '''keep all our PPT colours together'''
    from pptx.dml.color import RGBColor

    C_RGB_WHITE = RGBColor(255, 255, 255)
    C_RGB_BLACK = RGBColor(0, 0, 0)
    C_RGB_GREY = RGBColor(173, 173, 173)
    C_RGB_LIGHT_GREY = RGBColor(90, 90, 90)

class HiltonCommonCall(CommonCall):     # TODO implement this as a mixin ?
    '''encapsulate all hilton-specific coding'''
    from download_hilton_regional_updates import RESULT_ID as result_id
    from reportcheck import M_TELEPHONE,M_OVERALL,M_OVERALL_INC_QUICKCHECK,M_QUICKCHECK,M_ELECTRONIC,M_EMAIL,M_RFP,G_TELEPHONE,E_TELEPHONE,B_TELEPHONE

    result_id_set = [(M_TELEPHONE, 15604, u'Make package seem attractive', u'p', u'Selling Skills', 10),
                 (M_TELEPHONE, 15166, u'Fully outlined features & benefits', u'q', None, 0),
                 (M_TELEPHONE, 15605, u'Mentioned seasonal offers/loyalty programme', u'r', None, 0),
                 (M_TELEPHONE, 15608, u'Suggested relevant upgrades', u's', None, 0),
                 (M_TELEPHONE, 15170, u'Proposed provisional booking', u't', None, 0),
                 (M_TELEPHONE, 15401, u'Asked if future business requirements', u'u', None, 0),
                 (M_TELEPHONE, 15754, u'Product Knowledge %Exceptional', u'v', None, 0),
                 (M_TELEPHONE, 15756, u'Interest in enquiry throughout %Exceptional', u'w', None, 0),
                 (M_TELEPHONE, 15758, u'Handled specific request %Exceptional', u'x', None, 0),
                 (M_TELEPHONE, 15250, u'Made follow-up call', u'y', None, 0),
                 (M_TELEPHONE, 15191, u'% Completed in a single call', u'a', u'Customer Ease', 7),
                 (M_TELEPHONE, 15751, u'Helpfulness %Exceptional', u'f', None, 0),
                 (M_TELEPHONE, 15752, u'Clarity of explanations %Exceptional', u'g', None, 0),
                 (M_TELEPHONE, 15755, u'Sufficient use of time %Exceptional', u'h', None, 0),
                 (M_TELEPHONE, 15390, u'Percent as e-proposal or single document',u'j', None, 0),
                 (M_TELEPHONE, 15122, u'Sign-off on follow-up email', u'i', None, 0),
                 (M_TELEPHONE, 15247, u'Reference to special request in proposal', u'k', None, 0)]

    require_font_color_white = ['International Average', 'EMEA', 'APAC', 'International', 'EMEA', 'Asia-Pacific']
    require_font_type_bold = ['International Average', 'EUR inc UK', 'EMEA', 'APAC', 'MEA', 'ASIA', 'JAP', 'AUS']

    questionnaire_titles = {    #Hilton specific
        M_OVERALL_INC_QUICKCHECK:   "C&E Overall",
        M_OVERALL:                  "C&E Overall (excl. Quick Check)",
        M_TELEPHONE:                "C&E Telephone",
        M_QUICKCHECK:               "C&E Short Query",
        M_EMAIL:                    "C&E Email",
        M_RFP:                      "C&E RFP",
        E_TELEPHONE:                "Events Telephone",
        G_TELEPHONE:                "Group Res Telephone",
        B_TELEPHONE:                "Individual Res Telephone",
    }

class GeneratePPT(HiltonCommonCall):
    no_of_extra_slide = {}

    def __init__(self, dbConnection, **kwargs):
        self.d = dbConnection
        self.questionnaire_id = int(kwargs['questionnaire_id'])
        self.result_id = int(kwargs.get('result_id',self.result_id))
        self.id_type = kwargs['id_type']
        self.brand_id = int(kwargs['brand_id'])
        self.answer_value = kwargs['answer_value']
        self.score_year = int(kwargs['score_year'])
        self.score_month = int(kwargs['score_month'])
        self.demense_id = int(kwargs['demense_id'])
        self.division_id = int(kwargs['division_id'])
        self.score_type = kwargs['score_type']
        self.last_year, self.last_month = common.previous_period(year=self.score_year,month=self.score_month)
        self.demense_name = kwargs['demense_name']
        self.score_month_value = kwargs['score_month_value']
        self.brand_name = kwargs['brand_name']
        self.grouping_level = kwargs['grouping_level']


        ## Entity Object and Entity Type based on Entity Type
        self.demense_entity_id = get_scores.Entities.Demense
        self.demense_entity_type = get_scores.Entities(key_id=self.demense_entity_id)
        self.venue_entity_id = get_scores.Entities.Venue
        self.venue_entity_type = get_scores.Entities(key_id=self.venue_entity_id)
        self.division_entity_id = get_scores.Entities.OldDivision
        self.division_entity_type = get_scores.Entities(key_id=self.division_entity_id)
        self.brand_entity_id = get_scores.Entities.Brand
        self.brand_entity_type = get_scores.Entities(key_id=self.brand_entity_id)

        ## Manage division entity type for division Id for division report 
        ##  else Brand entity type for demense and Brand entity Type
        self.entity_type = self.brand_entity_type
        self.entity_id = self.brand_id
        self.did = demense_id

        if self.division_id:
            if self.division_id < 0:
                self.did = -self.division_id
                self.entity_type = self.division_entity_type
                self.entity_id = self.did
            else:
                self.did = self.division_id

        self.entity_id =  int(self.entity_id)
        self.did = int(self.did)
        self.param_arg = {'year': score_year,
                          'month': score_month,
                          'result_id': self.result_id,
                          'answer_value': self.answer_value,
                          'id_type': self.id_type,
                          'allow_override': True,
                          'own_scores': True}

        self.last_month_param_arg = {'year': self.last_year,
                                     'month': self.last_month,
                                     'result_id': self.result_id,
                                     'answer_value': self.answer_value,
                                     'id_type': self.id_type,
                                     'allow_override': True,
                                     'own_scores': True}

        self.brand_param_arg = {'entity_type': self.brand_entity_type,
                                'entity_id': self.brand_id}
        self.brand_param_arg.update(self.param_arg)

        self.demense_param_arg = {'entity_type': self.demense_entity_type,
                                  'entity_id': self.demense_id}
        self.demense_param_arg.update(self.param_arg)


    def getVenueScore(self, **kwargs):
        """
        This is common method to get VenueScores instead of calling many places

        this has been moved out of the superclass to here because it uses member vars that only exist at this level
        """
        questionnaire_id = kwargs.get('questionnaire_id',self.questionnaire_id)
        demense_id = kwargs.get('demense_id',self.demense_id)
        limit = kwargs.get('limit', 10)
        month_order_by = kwargs.get('month_order_by')
        ytd_order_by = kwargs.get('ytd_order_by')

        venuescores = model.VenueScores(questionnaire_id=questionnaire_id,
                                        result_id=self.result_id,
                                        id_type=self.id_type,
                                        brand_id=self.brand_id,
                                        answer_value=self.answer_value,
                                        set_mode=True).top_N_venue_in_brand(year=self.score_year,
                                                                            month=self.score_month,
                                                                            demense_id=demense_id,
                                                                            score_type=self.score_type,
                                                                            brand_id=self.brand_id,
                                                                            limit=limit,
                                                                            month_order_by=month_order_by,
                                                                            ytd_order_by=ytd_order_by) # FIXME 
        return venuescores

    def get_own_brand_demense_scores(self,
                                     did,
                                     questionnaire_id,
                                     score_type=None,
                                     year=None,
                                     month=None):
        """
        Get Own Brand Demense Score API
        """
        if not score_type:
            score_type = self.score_type

        if not year:
            year = score_year

        if not month:
            month = score_month

        score = model.Own_Brand_Demense_scores(id=did,
                                               year=year,
                                               month=month,
                                               score_type=score_type,
                                               result_id = self.result_id,
                                               questionnaire_id=questionnaire_id ,
                                               answer_value = self.answer_value,
                                               id_type = self.id_type,
                                               brand_id = brand_id,
                                               allow_override = True,
                                               own_scores = True).score or 0 # FIXME 
        return score

    def get_named_shape_obj(self,parent_obj,name_list):
        '''return a dict of named objects from the parent shape list'''
        return dict([(x.name,x) for x in parent_obj if x.name in name_list])

    def _set_tbl_row_font_size(self,
                              tableshape,
                              rowmapped,
                              columnmapped):
        ''' Set font size for corrosponding cell'''
        for i in columnmapped:
            tableshape.table.cell(rowmapped, i).text_frame.paragraphs[0].font.size = self.font_size

    def _clear_tbl_row(self,
                      tableshape,
                      rowmapped,
                      columnmapped):
        ''' Clear row for corrosponding cell'''
        for i in columnmapped:
            tableshape.table.cell(rowmapped, i).text_frame.clear()

    #-------------------------------------------------------------------------------------------------
    def mappedSlide(self, source_index):
        """
        Signature slide for 1 and 12
        This slide will show always
        """
        slideobj = prs.slides[source_index]
        return outputFile

    def mappedSlide1(self, source_index):
        """
        Show Demense name,month and year for slide 1
        """
        return self.mappedSlide(source_index)

    #-------------------------------------------------------------------------------------------------
    def mappedSlide2(self, source_index):
        """
        Show Demense name and scores for slide 2
        """

        TAB1 = 'BDRC_1_tbl_1'
        TAB2 = 'BDRC_1_tbl_2'
        IA_ICON = 'BDRC_1_International_Average'

        circular_icons = [
            IA_ICON,
            'BDRC_1_rgn_UK&I',
            'BDRC_1_rgn_W+S',
            'BDRC_1_rgn_MEA',
            'BDRC_1_rgn_C&EE',
            'BDRC_1_rgn_IND',
            'BDRC_1_rgn_GC',
            'BDRC_1_rgn_EUR+UK',
            'BDRC_1_rgn_APAC',
            'BDRC_1_rgn_SEA',
            'BDRC_1_rgn_JAP',
            'BDRC_1_rgn_AUS',
            'BDRC_1_rgn_EMEA',
            'BDRC_1_rgn_A+IO',
            'BDRC_1_rgn_KSA',
            'BDRC_1_rgn_ARAB',
            'BDRC_1_rgn_E+L',
            'BDRC_1_rgn_ASIA',
            'BDRC_1_rgn_RUKI',
            'BDRC_1_rgn_LON',
            'BDRC_1_rgn_AX',
        ]
        slideobj = prs.slides[source_index]

        ## Questionnaire ID is Overall for Top Worlwide Hotels
        ## Explicity set questionnari ID
        questionnaire_id = self.M_OVERALL
        ## TOP Worldwide Hotels Month and YTD based
        ##  Demense ID = 1 

        ## For MAP circular icons brand score value
        ##  minus is division and plus is demense
        for shape in slideobj.shapes:
            if shape.name in circular_icons and isinstance(shape,pptx.shapes.autoshape.Shape):
                did = int(shape.text.split('#')[1]) if shape.text.find( '#') != -1 else 0
                demense_name = shape.text.split( '#')[0] if shape.text.find('#') != -1 else 0
                if did:
                    if did < 0:
                        score = model.Division_scores(id=-did, questionnaire_id=questionnaire_id, score_type=self.score_type, **self.param_arg).score # FIXME 
                    else:
                        score = self.get_own_brand_demense_scores(did=int(did), questionnaire_id=questionnaire_id)
                    shape.text = demense_name + ' ' + str("%.1f" % float(score))
                    shape.text_frame.paragraphs[0].font.size = self.font_size_large if shape.name==IA_ICON else self.font_size
                    shape.text_frame.paragraphs[0].font.color.rgb = Colours.C_RGB_WHITE if [
                        item for item in self.require_font_color_white if demense_name.strip().find(item) != -1] else Colours.C_RGB_BLACK
                    if [item for item in self.require_font_type_bold if demense_name.strip().find(item) != -1]:
                        shape.text_frame.paragraphs[0].font.bold = True

        ## Set table record for TOP 10 Venues score 
        ## for month and YTD
        tables = self.get_named_shape_obj(slideobj.shapes,[TAB1,TAB2])
        if TAB1 in tables:
            monthvenuescores = self.getVenueScore(month_order_by=1, demense_id=1, questionnaire_id=questionnaire_id)
            self._mappedSlide_insert_table_vals(tables[TAB1],monthvenuescores)
        else:
            pass    # FIXME .... Raise an error ???

        if TAB2 in tables:
            ytdvenuescores = self.getVenueScore(ytd_order_by=1, demense_id=1, questionnaire_id=questionnaire_id)
            self._mappedSlide_insert_table_vals(tables[TAB2],ytdvenuescores)
        else:
            pass    # FIXME .... Raise an error ???

        prs.save(outputFile)
        return outputFile

    def _mappedSlide_insert_table_vals(self,tbl_obj,scores):
        row = 0
        for i in range(len(scores), len(tbl_obj.table.rows)):
            tbl_obj.table.cell(i, 0).text_frame.clear()
            tbl_obj.table.cell(i, 1).text_frame.clear()

        for each_val in scores:
            first_text_frame = tbl_obj.table.cell(row, 0).text_frame
            first_text_frame.paragraphs[0].runs[0].text = self.show_truncate_value(each_val[1])

            sec_text_frame = tbl_obj.table.cell(row, 1).text_frame
            sec_text_frame.paragraphs[0].runs[0].text = ("%.1f" % float(each_val[2]))

            row += 1


    #-------------------------------------------------------------------------------------------------
    def mappedSlide3(self, source_index):
        """
        Show Demense,questionnarie  and scores for slide 3
        """
        slideobj = prs.slides[source_index]

        TAB = 'BDRC_2_tbl'

        ## Get Demense Map from row 0 and based on that each score will
        ## calculate and set
        demense_map = {}
        q = 0
        for shape in slideobj.shapes:
            if not (shape.name==TAB and shape.has_table):
                continue
            for j in range(2, len(shape.table.columns)):
                colmap = shape.table.cell(0, j).text_frame.paragraphs[0].text
                colmap = colmap.split('=')
                entity_type = colmap[0]
                entity_type = entity_type.strip()
                if entity_type == DEM:
                    cdemense_id = colmap[1]
                    columnname = Demense(None).get(cdemense_id)[1]
                    demense_map[j] = [entity_type, cdemense_id, columnname]
                elif entity_type == DIV:
                    division_id = colmap[1]
                    columnname = Division(self.d, None).get(division_id)[1]
                    demense_map[j] = [entity_type, division_id, columnname]
            for i in range(1, len(shape.table.rows)):
                if i > 1 and i < 11:
                    score_type = '1'
                elif i > 11:
                    score_type = 'Y'
                else:
                    continue

                rowname = shape.table.cell(i, 0).text_frame.paragraphs[0].text
                if rowname.find('=') != -1:
                    colmap = rowname.split('=')
                    entity_type = str(colmap[0].strip())
                    entity_id = int(colmap[1])
                    if entity_type == QID:
                        questionnaire_id = entity_id

                shape.table.cell(i, 0).text_frame.paragraphs[0].text = self.questionnaire_titles[questionnaire_id]
                shape.table.cell(i, 0).text_frame.paragraphs[0].font.size = self.font_size
                q = q+1

                ## Based on input type ID and header will set
                ##  Calculate score as per desired entity Type and Entity ID
                for j in range(2, len(shape.table.columns)):
                    entity_type = demense_map[j][0]
                    if entity_type == DEM:
                        cdemense_id = demense_map[j][1]
                        columnname = demense_map[j][2]
                        score = self.get_own_brand_demense_scores(did=int(cdemense_id), questionnaire_id=questionnaire_id)
                    elif entity_type == DIV:
                        division_id = demense_map[j][1]
                        columnname = demense_map[j][2] 
                        score = model.Division_scores(id=division_id, questionnaire_id=questionnaire_id, score_type=score_type, **self.param_arg).score # FIXME
                    shape.table.cell(0, j).text_frame.paragraphs[0].text = columnname
                    shape.table.cell(0, j).text_frame.paragraphs[0].font.size = self.font_size
                    ## If labels are bottom look and feel is not correct for lengthy text 
                    ##  and not supporttaive to keep little space as prefix.
                    shape.table.cell(0, j).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER  
                    shape.table.cell(0, j).text_frame.paragraphs[0].font.bold = True
                    shape.table.cell(0, j).text_frame.paragraphs[0].font.color.rgb = Colours.C_RGB_WHITE if [
                        item for item in self.require_font_color_white if columnname.strip().find(item) != -1] else Colours.C_RGB_BLACK

                    cellObj = shape.table.cell(i, j).text_frame.paragraphs
                    score = str(int(round(score))) if score else 'n/a'
                    cellObj[0].text = score
                    cellObj[0].font.size = self.font_size
                    cellObj[0].runs[0].font.fill.solid()
                    cellObj[0].runs[
                        0].font.color.rgb = Colours.C_RGB_GREY if score == 'n/a' else Colours.C_RGB_BLACK
                    ## Anything N/A should be white cell background and grey text
                    if score == 'n/a':
                        shape.table.cell(
                            i, j).fill.fore_color.rgb = Colours.C_RGB_WHITE
        return outputFile

    #-------------------------------------------------------------------------------------------------
    def mappedSlide4(self, source_index):
        """
        Show Demense,questionnarie  and scores for slide 4
        """
        slideobj = prs.slides[source_index]
        TAB = 'BDRC_3_tbl'

        ## Get Demense Map from row 0.
        ## Based on that each difference score will calculate for last month - score month
        ## and set for each section
        demense_map = {}
        for shape in slideobj.shapes:
            if not (shape.name==TAB and shape.has_table):
                continue

            # process column headers
            for j in range(2, len(shape.table.columns)):
                textvalue = shape.table.cell(
                    0, j).text_frame.paragraphs[0].text
                if textvalue.find('=') == -1:
                    continue

                colmap = textvalue.split('=')
                entity_type = str(colmap[0].strip())
                entity_id = int(colmap[1])
                if entity_type == DEM:
                    cdemense_id = colmap[1]
                    columnname = Demense(None).get(cdemense_id)[1]
                    demense_map[j] = [entity_type, cdemense_id, columnname]
                elif entity_type == DIV:
                    division_id = colmap[1]
                    columnname = Division(self.d, None).get(division_id)[1]
                    demense_map[j] = [entity_type, division_id, columnname]

            # process row headers
            q = 0
            for i in range(1, len(shape.table.rows)):
                rowname = shape.table.cell(i, 0).text_frame.paragraphs[0].text
                ignore_row = 11
                if rowname == 'YTD <YYYY>':
                    ignore_row = i
                if i > 1 and i < ignore_row:
                    score_type = '1'
                elif i > ignore_row:
                    score_type = 'Y'
                else:
                    continue

                if rowname.find('=') != -1:
                    colmap = rowname.split('=')
                    entity_type = str(colmap[0].strip())
                    entity_id = int(colmap[1])
                    if entity_type == QID:
                        questionnaire_id = entity_id

                shape.table.cell(i, 0).text_frame.paragraphs[0].text = self.questionnaire_titles[questionnaire_id]
                shape.table.cell(i, 0).text_frame.paragraphs[0].font.size = self.font_size
                q = q+1

                ## Based on input type ID and header will set
                ##  Calculate score as per desired entity Type and Entity ID
                for j in range(2, len(shape.table.columns)):
                    entity_type = demense_map[j][0]
                    if entity_type == DEM:
                        cdemense_id = demense_map[j][1]
                        columnname = demense_map[j][2]
                        b_month_score = self.get_own_brand_demense_scores(did=cdemense_id, questionnaire_id=questionnaire_id, score_type=score_type)
                        b_last_month_score = self.get_own_brand_demense_scores(did=cdemense_id, questionnaire_id=questionnaire_id, score_type=score_type, year = self.last_year, month = self.last_month)
                        score = b_month_score - b_last_month_score
                    elif entity_type == DIV:
                        division_id = demense_map[j][1]
                        columnname = demense_map[j][2]
                        d_month_score = model.Division_scores(id=division_id, questionnaire_id=questionnaire_id, score_type=score_type, **self.param_arg).score or 0 # FIXME 
                        d_last_month_score = model.Division_scores(id=division_id, questionnaire_id=questionnaire_id, score_type=score_type, **self.last_month_param_arg).score or 0 # FIXME
                        score = d_month_score - d_last_month_score

                    shape.table.cell(0, j).text_frame.paragraphs[0].text = columnname
                    shape.table.cell(0, j).text_frame.paragraphs[0].font.size = self.font_size
                    ## If labels are bottom look and feel is not correct for lengthy text 
                    ##  and not supporttaive to keep little space as prefix.
                    shape.table.cell(0, j).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                    shape.table.cell(0, j).text_frame.paragraphs[0].font.bold = True
                    shape.table.cell(0, j).text_frame.paragraphs[0].font.color.rgb = Colours.C_RGB_WHITE if [item for item in self.require_font_color_white if columnname.strip().find(item) != -1] else Colours.C_RGB_BLACK
                    cellObj = shape.table.cell(i, j).text_frame.paragraphs
                    if not score:
                        score = 'n/a'
                    else:
                        score = str(int(round(score)))

                    cellObj[0].text = score
                    cellObj[0].font.size = self.font_size
                    cellObj[0].runs[0].font.fill.solid()
                    cellObj[0].runs[
                        0].font.color.rgb = Colours.C_RGB_GREY if score == 'n/a' else Colours.C_RGB_BLACK
                    ## Anything N/A should be white cell background and grey text
                    if score == 'n/a':
                        shape.table.cell( i, j).fill.fore_color.rgb = Colours.C_RGB_WHITE
        return outputFile

    #-------------------------------------------------------------------------------------------------
    def mappedSlide5(self, source_index):
        """
        Show Demense,questionnarie  and scores for slide 5
        """
        questionnaire_id = self.M_OVERALL
        TAB1 = 'BDRC_4_tbl_1'
        TAB2 = 'BDRC_4_tbl_2'

        slideobj = prs.slides[source_index]

        param_arg = {'demense_id': self.demense_id}
        param_arg.update(self.brand_param_arg)
        param_arg.pop('result_id')

        previous_param_arg = dict(param_arg)
        previous_param_arg.update({'year': self.last_year, 'month': self.last_month})
  
        ## Prcess score_method to generate score for each shape with required Input
        score_method = {'BDRC_4_Overall_Score': get_scores.get_score(questionnaire_id=self.M_OVERALL_INC_QUICKCHECK, result_id=self.result_id, score_type=self.score_type, **param_arg),
                        'BDRC_4_YTD': get_scores.get_score(questionnaire_id=self.M_OVERALL_INC_QUICKCHECK, result_id=self.result_id, score_type='Y', **param_arg),
                        'BDRC_4_Overall_EXCL': get_scores.get_score(questionnaire_id=questionnaire_id, result_id=self.result_id, score_type=self.score_type, **param_arg),
                        'BDRC_4_Overall_EXCL_Prev_Month': get_scores.get_score(questionnaire_id=questionnaire_id, result_id=self.result_id, score_type=self.score_type, **previous_param_arg),
                        'BDRC_4_Overall_EXCL_YTD': get_scores.get_score(questionnaire_id=questionnaire_id, result_id=self.result_id, score_type='Y', **param_arg),
                        'BDRC_4_Score_Telephone': get_scores.get_score(questionnaire_id=self.M_TELEPHONE, result_id=self.result_id, score_type='1', **param_arg),
                        'BDRC_4_Score_Overall': get_scores.get_score(questionnaire_id=self.M_QUICKCHECK, result_id=self.result_id, score_type='1', **param_arg),
                        'BDRC_4_Score_Email': get_scores.get_score(questionnaire_id=self.M_EMAIL, result_id=self.result_id, score_type='1', **param_arg),
                        'BDRC_4_Score_Web': get_scores.get_score(questionnaire_id=self.M_RFP, result_id=self.result_id, score_type='1', **param_arg),
                        'BDRC_4_Selling_Skills': get_scores.get_score(questionnaire_id=self.M_TELEPHONE, result_id='19025', score_type=self.score_type, **param_arg),
                        'BDRC_4_Customer_Ease': get_scores.get_score(questionnaire_id=self.M_TELEPHONE, result_id='19026', score_type=self.score_type, **param_arg),
                        'BDRC_4_Connection': get_scores.get_score(questionnaire_id=self.M_TELEPHONE, result_id='19021', score_type=self.score_type, **param_arg),
                        'BDRC_4_Service_Delivery': get_scores.get_score(questionnaire_id=self.M_TELEPHONE, result_id='19022', score_type=self.score_type, **param_arg),
                        'BDRC_4_Manner_Approach': get_scores.get_score(questionnaire_id=self.M_TELEPHONE, result_id='19023', score_type=self.score_type, **param_arg),
                        'BDRC_4_Followup': get_scores.get_score(questionnaire_id=self.M_TELEPHONE, result_id='19024', score_type=self.score_type, **param_arg),
#                        'BDRC_4_Overall_Prev_Month': ????
                        }

        ## For Sectional background int has to set so it's defined.
        didgit_required_list = ['BDRC_4_Connection', 'BDRC_4_Customer_Ease',
                                'BDRC_4_Selling_Skills', 'BDRC_4_Service_Delivery', 'BDRC_4_Manner_Approach', 'BDRC_4_Followup']

        for shape in slideobj.shapes:
            ## if shape.has_text_frame:
                ## for paragraph in shape.text_frame.paragraphs:
                    ## print shape.name
                    ## ,paragraph.text,score_method.get(shape.name)
            if shape.name in score_method.keys():
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        scorevalue = score_method[shape.name]
                        if shape.name == 'BDRC_4_Overall_EXCL':
                            continue

                        scorevalue = scorevalue[0] if scorevalue else 0.0
                        if shape.name == 'BDRC_4_Overall_EXCL_Prev_Month' and scorevalue:
                            newscorevalue = get_scores.get_score( questionnaire_id=self.M_OVERALL, result_id=self.result_id, score_type=self.score_type, **param_arg)
                            prevscorevalue = scorevalue
                            if prevscorevalue and newscorevalue:
                                scorevalue = newscorevalue[0] - prevscorevalue

                        elif shape.name == 'BDRC_4_Prev_Month':
                            newscorevalue = get_scores.get_score( questionnaire_id=self.M_OVERALL_INC_QUICKCHECK, result_id=self.result_id, score_type=self.score_type, **param_arg)
                            prevscorevalue = get_scores.get_score( questionnaire_id=self.M_OVERALL_INC_QUICKCHECK, result_id=self.result_id, score_type=self.score_type, **previous_param_arg)
                            if prevscorevalue and scorevalue:
                                scorevalue = newscorevalue[ 0] - prevscorevalue[0]
                        ## Calculate difference between Prev score and current
                        ## Manage Float and Integer format as per shape name
                        if shape.name in didgit_required_list and scorevalue:
                            scorevalue = str(int(round(scorevalue))) + '%'
                        elif (shape.name == 'BDRC_4_Prev_Month') or (shape.name == 'BDRC_4_Overall_EXCL_Prev_Month'):
                            if scorevalue > 0:
                                scorevalue = '+' + str(("%.1f" % float(scorevalue))) + '%'
                            else:
                                scorevalue = str( ("%.1f" % float(scorevalue))) + '%'

                        elif scorevalue:
                            scorevalue = str( ("%.1f" % float(scorevalue))) + '%'
                        else:
                            scorevalue = str(scorevalue)

                        paragraph.runs[0].text = scorevalue

        ## Set Table Values with Month scores  
        tables = self.get_named_shape_obj(slideobj.shapes,[TAB1,TAB2])
        if TAB1 in tables:
            monthvenuescores = self.getVenueScore(month_order_by=1, questionnaire_id=questionnaire_id, score_type=self.score_type)
            self._mappedSlide_insert_table_vals(tables[TAB1],monthvenuescores)
        else:
            pass    # FIXME .... Raise an error ???

        ## Set Table Values with YTD scores
        if TAB2 in tables:
            ytdvenuescores = self.getVenueScore(ytd_order_by=1, questionnaire_id=questionnaire_id, score_type='Y')
            self._mappedSlide_insert_table_vals(tables[TAB2],ytdvenuescores)
        else:
            pass    # FIXME .... Raise an error ???

        return outputFile

    #-------------------------------------------------------------------------------------------------
    def mappedSlide6(self, source_index):
        """
        Show Demense name and scores for slide 6
        """
        TAB1 = 'BDRC_5_tbl_1'
        TAB2 = 'BDRC_5_tbl_2'
        split50 = []
        split80 = []
        no_of_required_slide = 0
        questionnaire_id = self.M_OVERALL_INC_QUICKCHECK

        slideobj = prs.slides[source_index]

        overall_venue_score = self.getVenueScore(month_order_by=1, limit=None)
        addn = [self.M_TELEPHONE,self.M_ELECTRONIC, self.M_QUICKCHECK]
        ## venue score is calculated for and based on range it will set for
        ## specific table
        for ven in overall_venue_score:
            if ven[2] >= 80:  ## ignore anyone with >80% score
                continue

            val = []
            for n in addn:
                score = get_scores.get_score(entity_type=self.venue_entity_type,
                                             entity_id=ven[0],
                                             year=score_year,
                                             month=score_month,
                                             score_type='1',
                                             result_id=self.result_id,
                                             questionnaire_id=n,
                                             answer_value=self.answer_value,
                                             id_type=self.id_type,
                                             demense_id=demense_id,
                                             allow_override=True,
                                             own_scores=True)
                score = ("%.1f" % float(score[0])) if score else 'n/a'
                val.append(score)

            if ven[2] >= 50:  ## split out the >50%
                split80.append((ven[0], ven[1], ("%.1f" % float(ven[2])), val[0], val[1], val[2]))
            else:  ## split out the 50%-80%
                split50.append((ven[0], ven[1], ("%.1f" % float(ven[2])), val[0], val[1], val[2]))

        ##  split80tbl is for "Between 50% and 80%" table
        ##  split50tbl is for "Less than 50%" table
        tables = self.get_named_shape_obj(slideobj.shapes,[TAB1,TAB2])
        if tables:
            sp80len,sp80tablen = self._pre_process_slide6_table(tables[TAB1],split80)
            sp50len,sp50tablen = self._pre_process_slide6_table(tables[TAB2],split50)
            no_of_required_slide = max( int(math.ceil(float(sp80len) / float(sp80tablen) ) ), \
                                    int(math.ceil(float(sp50len) / float(sp50tablen) ) ) )
            ## Manage to set extra nos of record with new slides
            for i in range(1, no_of_required_slide):
                slide = self.dup_slide(prs,prs.slides[source_index])
                self._process_slide6_table(TAB1,slide,split80,split80tblshapelength,)
                self._process_slide6_table(TAB2,slide,split50,split50tblshapelength,)
                ## Manage new slides as per no of records
                ## Manage to set table with records for other new pages

        self.no_of_extra_slide[source_index] = no_of_required_slide
        return self.no_of_extra_slide

    def _process_slide6_table(self,TAB,slide,split,tblshapelength,):
        '''insert data (segment of split) into a single slide'''
        table_min_range = tblshapelength * (i - 1) + tblshapelength
        table_max_range = tblshapelength * i + tblshapelength
        tbl = split[table_min_range:table_max_range]
        tbllength = len(split)
        tblrangestart = tbllength + 1 if tbllength > 0 else tbllength
        noofrowclear = range(tblrangestart, tblshapelength + 1)
        tblrange = range(tbllength, tblshapelength)
        self.duplicate_data_slide6(slide,
                                   splitTBL=tbl,
                                   splitTBLnoofrowclear=noofrowclear,
                                   splitTBLrange=tblrange,
                                   TAB=TAB)

    def _pre_process_slide6_table(self, tab,split):
        '''work out how many sheets we need to lay out the data in split'''
        if tab:
            tblshapelength = len(tab.table.rows) - 1
            tbl = split[0:tblshapelength]
            tbllength = len(tbl)
            tblrangestart = tbllength + 1 if tbllength > 0 else tbllength + 2
            noofrowclear = range(tblrangestart, tblshapelength + 1)
            ## Manage to set table with records for 1st page
            self._duplicate_data_slide6(tab,                                # FIXME do we need to write the data into the tables twice ??
                                        splitTBL=tbl,
                                        splitTBLnoofrowclear=noofrowclear,
                                       )
            return tbllength,tblshapelength
        else:
            return None,None

    def duplicate_data_slide6(self,
                              newprs,
                              splitTBL,
                              splitTBLnoofrowclear=None,
                              splitTBLrange=None,
                              TAB=None):
        """
        Show Data for increament slide 6
        """
        ## Manage Table Object dynamically
        tables = self.get_named_shape_obj(newprs.shapes,[TAB])

        if tables:
            self._duplicate_data_slide6(tables[TAB],spliTBL,splitTBLnoofrowclear,splitTBLrange)

        return newprs


    def _duplicate_data_slide6(self,splitTBLshape,
                              splitTBL,
                              splitTBLnoofrowclear=None,
                              splitTBLrange=None,):
        row = 1
        for each_val in splitTBL:
            field = 1
                ## Manage table row values dymanically for corrponding table
            for frame in range(5):
                text_frame = splitTBLshape.table.cell( row, frame).text_frame
                text_frame.paragraphs[0].text = self.show_truncate_value( str(each_val[field]))
                text_frame.paragraphs[0].runs[ 0].font.size = self.font_size
                field += 1

            row += 1
            
            ## Manage table row to clear
        if splitTBLrange:
            for i in splitTBLrange:
                i = i + 1
                for j in range(5):
                    splitTBLshape.table.cell(i, j).text_frame.clear()

            ## Manage table row to keep white shadded if no record
        self.clearTablerow(splitTBLshape, splitTBLnoofrowclear, 5)

    #-------------------------------------------------------------------------------------------------
    def mappedSlide7(self, source_index):
        """
        Show Demense name and scores for slide 7
        """
        ## Set score range Object for Graph
        score_type = '3'
        overall_venue_score_range = []
        selling_skill_score_range = []
        customer_ease_score_range = []
        telephone_score_range = []
        electronic_score_range = []
        overall_demense_score_range = []
        selling_skill_demense_score_range = []
        customer_ease_demense_score_range = []
        telephone_demense_score_range = []
        electronic_demense_score_range = []

        slideobj = prs.slides[source_index]

        series1_name = str(brand_name) + ' in  ' + str(demense_name)  ## <BRAND> in <DEMENSE>
        series2_name = str(demense_name) + ' Average'  ## <DEMENSE> Average

        ## Manage to draw Graph for each section 
        ##  Overall Graph is brand Vs demense for period of 
        ##  1 year with 3 month interval
        ##  Others are for 6 month
        info = [(self.M_OVERALL_INC_QUICKCHECK, 19000, u"Overall Performance", overall_venue_score_range, overall_demense_score_range),
                (self.M_TELEPHONE, 19025, u"Selling Skills",selling_skill_score_range, selling_skill_demense_score_range),
                (self.M_TELEPHONE, 19026, u"Customer Ease",customer_ease_score_range, customer_ease_demense_score_range),
                (self.M_TELEPHONE, 19000, u"Telephone",telephone_score_range, telephone_demense_score_range),
                (self.M_ELECTRONIC, 19000, u"Electronic", electronic_score_range, electronic_demense_score_range)]

        brand_param_arg = {'score_type': '3'}
        brand_param_arg.update(self.brand_param_arg)
        brand_param_arg.pop('result_id')
        brand_param_arg.pop('year')
        brand_param_arg.pop('month')

        demense_param_arg = {'score_type': '3'}
        demense_param_arg.update(self.demense_param_arg)
        demense_param_arg.pop('result_id')
        demense_param_arg.pop('year')
        demense_param_arg.pop('month')

        ## This is used to calculate bar chart score for Brand Vs Demense
        for inf in info:
            for i in range(12):
                year,month = common.previous_period(year=score_year,month=score_month,size=i)
                scores = get_scores.get_score(year=year, month=month, demense_id=self.demense_id, result_id=inf[1], questionnaire_id=inf[0], **brand_param_arg)
                natscores = get_scores.get_score(year=year, month=month, result_id=inf[1], questionnaire_id=inf[0], **demense_param_arg)
                brand_score = ("%.1f" % float(scores[0])) if scores else ''
                demense_score = ("%.1f" % float( natscores[0])) if natscores else ''
                inf[3].append(brand_score)
                inf[4].append(demense_score)

            inf[3].reverse()
            inf[4].reverse()

        ## barchart_base is used to generate month description for X axis
        import barchart_base
        bb = barchart_base.Barchart_base()
        month_catagory = bb.set_bar_labels(limit=12, qurt_required=0)
        month_catagory = [i.replace('\n', '-') for i in month_catagory]
        ## Except Overall rest all are for last 6 month
        six_month_catagory = bb.set_bar_labels(limit=6, qurt_required=0)
        six_month_catagory = [i.replace('\n', '-') for i in six_month_catagory]
        
        brand_in_demense = ChartData()
        brand_in_demense.categories = month_catagory
        brand_in_demense.add_series(series1_name, overall_venue_score_range)
        brand_in_demense.add_series(series2_name, overall_demense_score_range)
        ## manage  selling skill score chart
        selling_skill_score_chart = ChartData()
        selling_skill_score_chart.categories = six_month_catagory
        selling_skill_score_chart.add_series('', selling_skill_score_range[5:])
        selling_skill_score_chart.add_series('', selling_skill_demense_score_range[5:])
        ## manage  customer ease score chart
        customer_ease_score_chart = ChartData()
        customer_ease_score_chart.categories = six_month_catagory
        customer_ease_score_chart.add_series('', customer_ease_score_range[5:])
        customer_ease_score_chart.add_series('', customer_ease_demense_score_range[5:])
        ## manage  telephone score chart
        telephone_score_chart = ChartData()
        telephone_score_chart.categories = six_month_catagory
        telephone_score_chart.add_series('', telephone_score_range[5:])
        telephone_score_chart.add_series('', telephone_demense_score_range[5:])
        ## manage  electronic score chart
        electronic_score_chart = ChartData()
        electronic_score_chart.categories = six_month_catagory
        electronic_score_chart.add_series('', electronic_score_range[5:])
        electronic_score_chart.add_series('', electronic_demense_score_range[5:])
        ## Manage to draw a graph with score values and generates shape
        for shape in slideobj.shapes:
            if not shape.has_chart:
                continue

            if shape.name == 'BDRC_6_chart_Overall':
                shape.chart.replace_data(brand_in_demense)
            elif shape.name == 'BDRC_6_chart_Selling_Skills':
                shape.chart.replace_data(selling_skill_score_chart)
            elif shape.name == 'BDRC_6_chart_Telephone':
                shape.chart.replace_data(telephone_score_chart)
            elif shape.name == 'BDRC_6_chart_Customer_Ease':
                shape.chart.replace_data(customer_ease_score_chart)
            elif shape.name == 'BDRC_6_chart_Electronic':
                shape.chart.replace_data(electronic_score_chart)

        return outputFile

    #-------------------------------------------------------------------------------------------------
    def mappedSlide8(self, source_index):
        """
        Show Demense name and scores for slide 8
        """
        catagorylist = []
        month_score_range = []
        ytd_score_range = []

        slideobj = prs.slides[source_index]

        param_arg = {'demense_id': self.demense_id}
        param_arg.update(self.brand_param_arg)
        param_arg.pop('result_id')

        previous_param_arg = dict(param_arg)
        previous_param_arg.update({'year': self.last_year, 'month': self.last_month})

        ## This is used to calculate chart score for each catagory 
        ##  for month and YTD difference
        for item in self.result_id_set:
            catagorylist.append(item[2])
            month_score = get_scores.get_score(score_type='1', result_id=item[1], questionnaire_id=item[0], **param_arg)
            last_month_score = get_scores.get_score(score_type='1', result_id=item[1], questionnaire_id=item[0], **previous_param_arg)
            month_score = month_score[0] if month_score else 0.0
            last_month_score = last_month_score[0] if last_month_score else 0.0
            if last_month_score != 0.0 and month_score != 0.0:
                month_diffScore = month_score - last_month_score
            else:
                month_diffScore  = 0.0

            month_score_range.append(month_diffScore)
            ytd_score = get_scores.get_score(score_type='Y', result_id=item[1], questionnaire_id=item[0], **param_arg)
            last_ytd_score = get_scores.get_score(score_type='Y', result_id=item[1], questionnaire_id=item[0], **previous_param_arg)
            ytd_score = ytd_score[0] if ytd_score else 0.0
            last_ytd_score = last_ytd_score[0] if last_ytd_score else 0.0
            if last_ytd_score != 0.0 and ytd_score != 0.0:
                ytd_diffScore = ytd_score - last_ytd_score
            else:
                ytd_diffScore  = 0.0

            ytd_score_range.append(ytd_diffScore)

        ## Month Score Chart for Selling Skills and Customer Ease
        month_score_chart = ChartData()
        month_score_chart.categories = catagorylist
        month_score_chart.add_series('', month_score_range)
        ## YTD Score Chart for Selling Skills and Customer Ease
        ytd_score_chart = ChartData()
        ytd_score_chart.categories = catagorylist
        ytd_score_chart.add_series('', ytd_score_range)

        ## Manage to set month and YTD record for Selling Skills and Customer Ease
        for shape in slideobj.shapes:
            if not shape.has_chart:
                continue

            if shape.name == 'BDRC_7_chart_Month':
                shape.chart.replace_data(month_score_chart)
            elif shape.name == 'BDRC_7_chart_YTD':
                shape.chart.replace_data(ytd_score_chart)

        return outputFile

    #-------------------------------------------------------------------------------------------------
    def mappedSlide9(self, source_index):
        """
        Show Demense name and scores for slide 9
        """
        catagorylist = []
        month_score_range = []
        ytd_score_range = []

        slideobj = prs.slides[source_index]

        param_arg = {'demense_id': self.demense_id}
        param_arg.update(self.brand_param_arg)
        param_arg.pop('result_id')

        demense_param_arg = dict(param_arg)
        demense_param_arg.update({'entity_type': self.demense_entity_type, 'entity_id': self.demense_id})

        ## This is used to calculate chart score for each catagory 
        ##  for Brand VS Demense difference
        for item in self.result_id_set:
            catagorylist.append(item[2])
            brand_score = get_scores.get_score(score_type='1', result_id=item[1], questionnaire_id=item[0], **param_arg)
            demense_score = get_scores.get_score(score_type='1', result_id=item[1], questionnaire_id=item[0], **demense_param_arg)
            brand_score = brand_score[0] if brand_score else 0.0
            demense_score = demense_score[0] if demense_score else 0.0
            month_diffScore = float(brand_score) - float(demense_score)
            month_score_range.append(month_diffScore)
            ytd_brand_score = get_scores.get_score(score_type='Y', result_id=item[1], questionnaire_id=item[0], **param_arg)
            ytd_demense_score = get_scores.get_score(score_type='Y', result_id=item[1], questionnaire_id=item[0], **demense_param_arg)
            ytd_brand_score = ytd_brand_score[0] if ytd_brand_score else 0.0
            ytd_demense_score = ytd_demense_score[0] if ytd_demense_score else 0.0
            ytd_diffScore = float(ytd_brand_score) - float(ytd_demense_score)
            ytd_score_range.append(ytd_diffScore)

        ## Month Score Chart for Selling Skills and Customer Ease
        month_score_chart = ChartData()
        month_score_chart.categories = catagorylist
        month_score_chart.add_series('', month_score_range)
        ## YTD Score Chart for Selling Skills and Customer Ease
        ytd_score_chart = ChartData()
        ytd_score_chart.categories = catagorylist
        ytd_score_chart.add_series('', ytd_score_range)
        ## Manage to set month and YTD record for Selling Skills and Customer Ease
        for shape in slideobj.shapes:
            if not shape.has_chart:
                continue

            if shape.name == 'BDRC_8_chart_Month':
                shape.chart.replace_data(month_score_chart)
            elif shape.name == 'BDRC_8_chart_YTD':
                shape.chart.replace_data(ytd_score_chart)

        return outputFile

    #-------------------------------------------------------------------------------------------------
    def mappedSlide10(self, source_index):
        """
        Show Demense name and scores for slide 10
        """
        allvenuescore = {}

        questionnaire_id = self.M_OVERALL

        slideobj = prs.slides[source_index]

        allbrands = brand.Brand(self.d, None).get_all_parent_brand()


        ## Manage table values in below order
        ##  1. Brand name in League Table
        ##  2. Venue last month score for meetings inc quick
        ##  3. Venue YTD score for meetings inc quick
        ##  4. Demense = 1 Rank
        ##  5. Venue Score for Telephone in last 3 month
        ##  7. Venue Score for Electronic in last 3 month

        for brand_id, brand_name in allbrands:
            score_args = {'entity_type': self.brand_entity_type,
                          'entity_id': brand_id,
                          'result_id': self.result_id,
                          'answer_value': self.answer_value,
                          'id_type': self.id_type,
                          'demense_id': self.demense_id
                          }
            rank_args = {'questionnaire_id': questionnaire_id,
                         'for_brand': brand_id,
                         'grouping_level': self.grouping_level,
                         'score_type': 'Y',
                         'year': score_year,
                         'month': score_month
                         }
            rank_args.update(score_args)
            rank_args.update({'entity_type': self.venue_entity_type, 'demense_id': 1})

            ## This is used to calculate score for last month overall score ,YTD
            ## , worldwide rank , last 3 month telephone score
            overall_last_month_score = get_scores.get_score(questionnaire_id=questionnaire_id, score_type='1', year=score_year, month=score_month, **score_args)
            overall_YTD_score = get_scores.get_score(questionnaire_id=questionnaire_id, score_type='Y', year=score_year, month=score_month, **score_args)
            worldwide_rank = get_scores.get_rank(**rank_args)[0]
            telephone_brandscore = get_scores.get_score(questionnaire_id=self.M_TELEPHONE, year=score_year, month=score_month, score_type='3', **score_args)
            allvenuescore[brand_name] = [overall_last_month_score,overall_YTD_score, worldwide_rank, telephone_brandscore]

        self.duplicate_data_slide10(prs.slides[source_index],
                                    score_month_value=score_month_value,
                                    score_year=score_year,
                                    allvenuescore=allvenuescore)
        return outputFile

    def duplicate_data_slide10(self,
                               newprs,
                               score_month_value,
                               score_year,
                               allvenuescore):
        """
        Duplicate data as per requirment
        """
        TAB = 'BDRC_9_tbl'
        tables = self.get_named_shape_obj(newprs.shapes,[TAB,])
        if tables:
            tableshape = tables[TAB]
            tableshape.table.cell(0, 0).text_frame.paragraphs[0].runs[ 0].text = str(score_month_value) + ' ' + str(score_year)
            self._set_tbl_row_font_size(tableshape=tableshape,
                                       rowmapped=0,
                                       columnmapped=[0,2,5,7])
            row = 2
            noOfRecords = len(allvenuescore)
            noOfRow = len(tableshape.table.rows)
            if noOfRecords:
                clearrowrange = range(noOfRecords, len(tableshape.table.rows))
                whiteshadednoofrows = noOfRecords + 2 if noOfRecords > 0 else noOfRecords
                whiteshadednoofrows = range(
                    whiteshadednoofrows, len(tableshape.table.rows))
                for i in clearrowrange:
                    self._clear_tbl_row(tableshape=tableshape,
                                       rowmapped=i,
                                       columnmapped=[0,2,5,7])

                self.clearTablerow(tableshape, whiteshadednoofrows, 8)
            ## Set Table values with Output
            for key, value in allvenuescore.iteritems():
                tableshape.table.cell(row, 0).text_frame.paragraphs[0].text = self.show_truncate_value(key)
                tableshape.table.cell(row, 2).text_frame.paragraphs[0].text = str( "%.1f" % float(value[0][0])) if value[0] else 'n/a'
                tableshape.table.cell(row, 3).text_frame.paragraphs[0].text = str( "%.1f" % float(value[1][0])) if value[1] else 'n/a'
                tableshape.table.cell(row, 5).text_frame.paragraphs[0].text = str( "%.1f" % float(value[2])) if value[2] else 'n/a'
                tableshape.table.cell(row, 7).text_frame.paragraphs[0].text = str( "%.1f" % float(value[3][0])) if value[3] else 'n/a'
                self._set_tbl_row_font_size(tableshape=tableshape,
                                           rowmapped=row,
                                           columnmapped=[0,2,3,5,7])
                row += 1
        return newprs


    #-------------------------------------------------------------------------------------------------
    def mappedSlide11(self, source_index):
        """
        Show Demense name and scores for slide 11
        @source_index : {Int} slide Index
        """
        allvenuescore = []

        slideobj = prs.slides[source_index]

        ## Manage Demense name as per user Dynamic user input
        TAB = 'BDRC_10_tbl'
        tables = self.get_named_shape_obj(slideobj.shapes,[TAB,])
        tableshape = tables[TAB]
        tableshape.table.cell(1, 5).text_frame.paragraphs[0].text = demense_name
        tableshape.table.cell(1, 5).text_frame.paragraphs[0].font.size = self.font_size
        tableshape.table.cell(1, 5).text_frame.paragraphs[0].font.color.rgb = Colours.C_RGB_WHITE

        ## Get all venue Score
        ovenuescores = self.getVenueScore(limit=None)

        ## Set Score keyword arguments
        score_args = {'entity_type': self.venue_entity_type,
                      'year': score_year,
                      'month': score_month,
                      'result_id': self.result_id,
                      'answer_value': self.answer_value,
                      'id_type': self.id_type,
                      }
        rank_args = {'questionnaire_id': self.M_OVERALL_INC_QUICKCHECK,
                     'for_brand': brand_id,
                     'grouping_level': self.grouping_level,
                     'score_type': 'Y'
                     }
        rank_args.update(score_args)

        ## Manage table values in below order
        ##  1. Venue name
        ##  2. Venue last month score for meetings inc quick
        ##  3. Venue YTD score for meetings inc quick
        ##  4. Demense = Demense_ID Rank
        ##  5. Demense = 1 Rank
        ##  6. Venue Score for Telephone in last 3 month
        ##  7. Venue Score for Electronic in last 3 month
        ##  8. Venue Score for Quickcheck in last 3 month
        for item in ovenuescores:
            ## This is used to calculate score for last month overall score ,YTD
            ## , demense rank ,worldwide rank , last 3 month
            ## telephone,electronic and quickcheck score
            demense_rank = get_scores.get_rank(entity_id=item[0], demense_id=demense_id, **rank_args)[0]
            worldwide_rank = get_scores.get_rank(entity_id=item[0], demense_id=1, **rank_args)[0]
            telephone_venuescores = get_scores.get_score(entity_id=item[0], questionnaire_id=self.M_TELEPHONE, score_type='3', **score_args)
            electronic_venuescores = get_scores.get_score(entity_id=item[0], questionnaire_id=self.M_ELECTRONIC, score_type='3', **score_args)
            quickcheck_venuescores = get_scores.get_score(entity_id=item[0], questionnaire_id=self.M_QUICKCHECK, score_type='3', **score_args)
            allvenuescore.append((item[1], item[2], item[3], demense_rank, worldwide_rank,telephone_venuescores, electronic_venuescores, quickcheck_venuescores))

        ## Manage to add set on 1st sldie as per record
        self.duplicate_data_slide11(prs.slides[source_index],
                                    score_month_value=score_month_value,
                                    score_year=score_year,
                                    allvenuescore=allvenuescore[0:18])

        ## Manage to add new slides as per record count per page
        no_of_required_slide = int(math.ceil(float(len(allvenuescore)) / float(len(tableshape.table.rows))))
        for i in range(1, no_of_required_slide):
            fromrange = 18 * (i - 1) + 18
            torange = 18 * i + 18
            newprs = self.duplicate_slide(prs, source_index)
            self.duplicate_data_slide11(newprs,
                                score_month_value=score_month_value,
                                score_year=score_year,
                                allvenuescore=allvenuescore[fromrange:torange])

        self.no_of_extra_slide[source_index] = no_of_required_slide
        return self.no_of_extra_slide

    def duplicate_data_slide11(self,
                       newprs,
                       score_month_value,
                       score_year,
                       allvenuescore):
        """
        Duplicate data as per requirment
        This Method is used to clear record and update data for ech cell in slide 11
        """
        TAB = 'BDRC_10_tbl'
        tables = self.get_named_shape_obj(newprs.shapes,[TAB,])
        if tables:
            tableshape = tables[TAB]
            tableshape.table.cell(0, 0).text_frame.paragraphs[0].runs[0].text = str(score_month_value) + ' ' + str(score_year)
            self._set_tbl_row_font_size(tableshape=tableshape,
                                       rowmapped=0,
                                       columnmapped=[0,2,5,8])

            row = 2
            noOfRecords = len(allvenuescore)
            noOfRow = len(tableshape.table.rows)
            if noOfRecords:
                clearrowrange = range(noOfRecords, len(tableshape.table.rows))
                whiteshadednoofrows = noOfRecords + 2 if noOfRecords > 0 else noOfRecords
                whiteshadednoofrows = range(whiteshadednoofrows, len(tableshape.table.rows))
                for i in clearrowrange:
                    self._clear_tbl_row(tableshape=tableshape,
                                       rowmapped=i,
                                       columnmapped=[0,2,3,5,6,8,9,10])

                self.clearTablerow(tableshape, whiteshadednoofrows, 11)

            ## Set Table values with Output
            for item in allvenuescore:
                tableshape.table.cell(row, 0).text_frame.paragraphs[0].text = self.show_truncate_value(item[0])
                tableshape.table.cell(row, 2).text_frame.paragraphs[0].text = str( "%.1f" % float(item[1])) if item[1] else 'n/a'
                tableshape.table.cell(row, 3).text_frame.paragraphs[0].text = str( "%.1f" % float(item[2])) if item[2] else 'n/a'
                tableshape.table.cell(row, 5).text_frame.paragraphs[0].text = str(item[3]) if item[3] else 'n/a'
                tableshape.table.cell(row, 6).text_frame.paragraphs[0].text = str(item[4]) if item[4] else 'n/a'
                tableshape.table.cell(row, 8).text_frame.paragraphs[0].text = str( "%.1f" % float(item[5][0])) if item[5] else 'n/a'
                tableshape.table.cell(row, 9).text_frame.paragraphs[0].text = str( "%.1f" % float(item[6][0])) if item[6] else 'n/a'
                tableshape.table.cell(row, 10).text_frame.paragraphs[0].text = str( "%.1f" % float(item[7][0])) if item[7] else 'n/a'
                self._set_tbl_row_font_size(tableshape=tableshape,
                                           rowmapped=row,
                                           columnmapped=[0,2,3,5,6,7,8,9,10])
                row += 1

        return newprs

    #-------------------------------------------------------------------------------------------------
    def mappedSlide12(self, source_index):
        """
        Signature slide 12
        This slide will show always
        """
        return self.mappedSlide(source_index)

#-------------------------------------------------------------------------------------------------
#-------------------------------------------------------------------------------------------------
#-------------------------------------------------------------------------------------------------
if __name__ == '__main__':
    ## cgitb module is used to check log
    import cgitb
    cgitb.enable()
    ## DB connection attribute
    d = db.connect()
    ## INput form field storage
    form = cgi.FieldStorage()
    ## Manage Brand,Demense,Month,year and Division from Input Form
    brand_id = form['brand_id'].value
    brands = brand.Brand(d, brand_id)
    brand_name = brands.brand_name
    if form.has_key('demense_id'):
        demense_id = form['demense_id'].value
        demense_name = Demense(None).get(demense_id)[1]
    else:
        demense_id = 0
        demense_name = ''

    score_month = int(form['month'].value)
    score_month_value = datetime.date(1900, score_month, 1).strftime('%B')
    score_year = int(form['year'].value)
    division_id = form['division_id'].value
    ## Manage slide order based based on User provided 
    if form.has_key('selected_index'):
        selected_index = form['selected_index'].value
    else:
        selected_index = None

    ## Manage to generate PPT
    ReportManager(d).generatePPT(brand_id,
                            brand_name,
                            demense_id,
                            demense_name,
                            score_month,
                            score_month_value,
                            score_year,
                            division_id,
                            selected_index)
    ## For slide order it rearrange slide index based on 
    ##  input and removes unwanted slides,else keep all
    ##  slide in descending order
    if selected_index:
        content = desttemp.file.read()
        ## content  = open(outputFile).read()
        ## filename = os.path.basename(desttempFile)
    else:
        content = temp.file.read()
        ## content  = open(outputFile).read()
        ## filename = os.path.basename(outputFile)

    ## Manage File name with .pptx extension and pptx header,
    ##  keep slide content and 
    ##  which allow to save file by user 
    downloadfilename = 'BDRC_2016_Output_' + str(brand_id) + '_' + str(demense_id) + '_' + datetime.datetime.now().strftime('%Y_%m') + '.pptx'
    sys.stdout.write('content-type: application/vnd.openxmlformats-officedocument.presentationml.presentation\n')
    sys.stdout.write('Content-Disposition: attachment; filename="%s"\n\n' % downloadfilename)
    sys.stdout.write(content)

    ## Destory all object
    d.close()
    temp.close()
    desttemp.close()

#-------------------------------------------------------------------------------------------------
#-------------------------------------------------------------------------------------------------
#-------------------------------------------------------------------------------------------------
