
import venue
import model
import copy
import datetime
#import constants
import get_scores
import collections
from pptx.util import Pt

class CommonCall():
    font_size = Pt(7.5)
    font_size_large = Pt(12.5)

    def show_truncate_value(self, value, max_char=25):
        """
        Truncate and return 25 char with ...
        """
        # if len(value) > max_char:
        #     value = value[:max_char] + '...'
        return value

    def _get_blank_slide_layout(self, prs):
        layout_items_count = [len(layout.placeholders) for layout in prs.slide_layouts]
        min_items = min(layout_items_count)
        blank_layout_id = layout_items_count.index(min_items)
        return prs.slide_layouts[blank_layout_id]

    def dup_slide(self,prs,slide):
        """clone the slide and append onto prs
            return the cloned slide"""
        blank_slide_layout = self._get_blank_slide_layout(prs)
        clone = prs.slides.add_slide(blank_slide_layout)
        for shp in slide.shapes:
            el = shp.element
            newel = copy.deepcopy(el)
            clone.shapes._spTree.insert_element_before(newel, 'p:extLst')

        for key, value in slide.rels.iteritems():
            # Make sure we don't copy a notesSlide relation as that won't exist
            if not "notesSlide" in value.reltype:
                clone.rels.add_relationship( value.reltype, value._target, value.rId)

        return clone


    def duplicate_slide(self, prs, index):
        """Duplicate the slide with the given index in prs.
        Adds slide to the end of the presentation
        return a pointer to the slide"""
        return self.dup_slide(prs, prs.slides[index])

    def clearTablerow(self, tableshape, noofrows, noofcolumns):
        """
        If row is empty (no data) to maintain not to have a row, or added white shading

        @tableshape  :  (obj) Table Object for which clear needed
        @noofrows : (list) List of rows 
        @noofcolumns : (int) No Of Columns used for range
        """
        from gen_pptx_report import Colours


        for i in noofrows:
            for j in range(noofcolumns):
                tableshape.table.cell( i, j).fill.fore_color.rgb = Colours.C_RGB_WHITE

